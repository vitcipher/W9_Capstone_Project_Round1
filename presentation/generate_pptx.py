"""
Generates slides.pptx from the same content as slides.html (the ledger/
property-record visual identity - IBM Plex Sans/Mono, deep-teal accent,
green/amber semantic badges), so there's a real, editable presentation
file alongside the standalone HTML deck. Uploading this .pptx to Google
Drive auto-converts it to a native, fully editable Google Slides file.

Run: python3 generate_pptx.py  (needs python-pptx: pip install python-pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette (light mode of the ledger theme in slides.html) ----
INK = RGBColor(0x1B, 0x2A, 0x2E)
PAPER = RGBColor(0xEC, 0xEF, 0xE2)
PAPER_RAISED = RGBColor(0xF5, 0xF6, 0xEC)
RULE = RGBColor(0xC7, 0xCB, 0xB8)
RULE_STRONG = RGBColor(0xA9, 0xAE, 0x98)
ACCENT = RGBColor(0x1F, 0x5E, 0x56)
CONFIRMED = RGBColor(0x3F, 0x7D, 0x4C)
FLAGGED = RGBColor(0xB5, 0x72, 0x2A)
MUTED = RGBColor(0x5B, 0x66, 0x5F)

DISPLAY_FONT = "IBM Plex Sans"
MONO_FONT = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs, bg=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def _set_run(run, text, size=14, color=INK, bold=False, italic=False, font=DISPLAY_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts with keys text,size,color,bold,italic,font,space_after"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line.get("space_after") is not None:
            p.space_after = Pt(line["space_after"])
        run = p.add_run()
        _set_run(
            run,
            line["text"],
            size=line.get("size", 14),
            color=line.get("color", INK),
            bold=line.get("bold", False),
            italic=line.get("italic", False),
            font=line.get("font", DISPLAY_FONT),
        )
    return tb


def eyebrow(slide, idx, label):
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.45), Inches(10), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    _set_run(r1, f" {idx} ", size=11, color=RGBColor(0xEC, 0xEF, 0xE2), bold=True, font=MONO_FONT)
    # background chip for the index number
    chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.42), Inches(0.42), Inches(0.28))
    chip.fill.solid()
    chip.fill.fore_color.rgb = ACCENT
    chip.line.fill.background()
    ctf = chip.text_frame
    ctf.margin_left = Emu(0)
    ctf.margin_right = Emu(0)
    ctf.margin_top = Emu(0)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    _set_run(cr, idx, size=10, color=PAPER, bold=True, font=MONO_FONT)
    tb.left = Inches(1.2)
    tf.paragraphs[0].runs[0].text = label
    _set_run(tf.paragraphs[0].runs[0], label, size=11, color=ACCENT, bold=True, font=MONO_FONT)
    return chip


def title(slide, text, top=Inches(0.85), size=34, width=Inches(10.5)):
    return textbox(slide, MARGIN, top, width, Inches(1.6), [
        {"text": text, "size": size, "bold": True, "color": INK}
    ])


def source_line(slide, text, top=Inches(6.95)):
    return textbox(slide, MARGIN, top, Inches(12), Inches(0.35), [
        {"text": text, "size": 9.5, "color": MUTED, "font": MONO_FONT}
    ])


def footer(slide, page, total, label="Property Ledger"):
    textbox(slide, MARGIN, Inches(7.1), Inches(4), Inches(0.3), [
        {"text": label, "size": 9.5, "color": MUTED, "font": MONO_FONT}
    ])
    textbox(slide, Inches(11.8), Inches(7.1), Inches(1.0), Inches(0.3), [
        {"text": f"{page:02d} / {total:02d}", "size": 9.5, "color": MUTED, "font": MONO_FONT}
    ], align=PP_ALIGN.RIGHT)


def field_box(slide, left, top, width, height, label, value, note="", value_color=INK, value_size=26):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PAPER_RAISED
    box.line.color.rgb = RULE_STRONG
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    _set_run(r0, label, size=10.5, color=MUTED, bold=False, font=MONO_FONT)
    p0.space_after = Pt(4)

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    _set_run(r1, value, size=value_size, color=value_color, bold=True, font=DISPLAY_FONT)
    p1.space_after = Pt(4)

    if note:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        _set_run(r2, note, size=10.5, color=INK, font=DISPLAY_FONT)
    return box


def badge(slide, left, top, text, kind="confirmed"):
    color = CONFIRMED if kind == "confirmed" else FLAGGED
    soft = RGBColor(0xDC, 0xEA, 0xD9) if kind == "confirmed" else RGBColor(0xF2, 0xE2, 0xCD)
    w, h = Inches(1.7), Inches(0.34)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = soft
    box.line.color.rgb = color
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, text.upper(), size=9.5, color=color, bold=True, font=MONO_FONT)
    return box


def bullet_list(slide, left, top, width, height, items, size=14, gap=10):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        _set_run(r, "—  " + item, size=size, color=INK, font=DISPLAY_FONT)
    return tb


def add_table(slide, left, top, width, height, header, rows, col_widths):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(header), left, top, width, height)
    table = tbl_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for c, val in enumerate(header):
        cell = table.cell(0, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER_RAISED
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = MONO_FONT
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = MUTED
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = MONO_FONT
            run.font.size = Pt(13)
            run.font.color.rgb = INK
    return tbl_shape


TOTAL = 15
prs = new_deck()

# ---------------------------------------------------------------- Slide 1
s = add_slide(prs)
textbox(s, MARGIN, Inches(0.5), Inches(6), Inches(0.4), [
    {"text": "ROUND 1 — AI CONSULTING PITCH", "size": 11, "color": MUTED, "font": MONO_FONT}
])
textbox(s, MARGIN, Inches(2.5), Inches(10.5), Inches(2.2), [
    {"text": "The Property Ledger Pitch", "size": 46, "bold": True, "color": INK}
])
textbox(s, MARGIN, Inches(4.15), Inches(9.5), Inches(1.2), [
    {"text": ("A property-management module for private landlords, pitched to Chleo — "
              "CEO of an existing property listing platform who isn't sure what "
              "“the AI” actually does."), "size": 17, "color": MUTED}
])
textbox(s, MARGIN, Inches(6.6), Inches(10), Inches(0.4), [
    {"text": "Prepared for the Round 1 teaching-staff presentation", "size": 11, "color": MUTED, "font": MONO_FONT}
])
footer(s, 1, TOTAL)

# ---------------------------------------------------------------- Slide 2
s = add_slide(prs)
eyebrow(s, "02", "THE SCENARIO")
title(s, "Chleo just finished asking the question every AI pitch dreads.")
textbox(s, MARGIN, Inches(1.9), Inches(10.8), Inches(1.1), [
    {"text": ("She runs a property listing platform. She's heard the AI hype, and she's "
              "scared it isn't transparent — she doesn't know what “the AI” is, "
              "or how she'd explain it to her own team if something went wrong."),
     "size": 17, "color": INK}
])
bullet_list(s, MARGIN, Inches(3.2), Inches(7.2), Inches(2), [
    "She isn't asking for a moonshot. She's asking for something she can look at and understand.",
    "This pitch is built to answer her fear directly, not talk around it — every claim is cited, every limit is stated out loud.",
])
field_box(s, Inches(8.5), Inches(3.2), Inches(4.1), Inches(1.9),
          "WHAT THIS PITCH IS NOT",
          "A black box that promises to handle everything",
          "It's one narrow, honest capability — see slide 05.", value_size=17)
footer(s, 2, TOTAL)

# ---------------------------------------------------------------- Slide 3
s = add_slide(prs)
eyebrow(s, "03", "WHY NOW")
title(s, "Private landlords in Germany aren't a shrinking niche.", size=30)
coords = [
    (MARGIN, Inches(2.0), "PRIVATE-LANDLORD HOUSEHOLDS, DE", "3.7M → 5.5M+", "10% → 13% of the population, 2010→2022 — a genuine structural shift."),
    (Inches(6.9), Inches(2.0), "GERMANY IS MAJORITY-RENTER", "58%", "of households rent (Zensus 2022) — the only EU country where more rent than own."),
    (MARGIN, Inches(3.75), "FOREIGN POPULATION, DE", "10.9M → 14.1M", "+28.9%, 2018→2025 — skews toward renting (Ukraine +897%, India +151%)."),
    (Inches(6.9), Inches(3.75), "REAL-ESTATE AI HIRING, 2025", "+93.5%", "YoY, 2nd-fastest of any sector tracked — still early-stage."),
]
for x, y, label, value, note in coords:
    field_box(s, x, y, Inches(6.0), Inches(1.55), label, value, note, value_color=ACCENT, value_size=28)
source_line(s, "SOURCE: IW Köln “Private Vermieter in Deutschland” (SOEP v39+Zensus) · Destatis Zensus 2022/AZR · Stanford HAI AI Index 2026", top=Inches(5.6))
footer(s, 3, TOTAL)

# ---------------------------------------------------------------- Slide 4
s = add_slide(prs)
eyebrow(s, "04", "WHO THE CUSTOMER REALLY IS")
title(s, "Not a professional investor — someone with one apartment and a spreadsheet they don't trust.", size=27)
field_box(s, MARGIN, Inches(2.15), Inches(5.6), Inches(1.7), "OWN AT MOST 2 PROPERTIES", "~75%",
          "of private landlords (IW Köln Vermieterreport 2026, n=1,002) — and the 1-property share is rising: 55%→58%, 2024→2026.",
          value_color=ACCENT, value_size=30)
field_box(s, Inches(7.1), Inches(2.15), Inches(5.6), Inches(1.7), "RENTAL INCOME IS MINOR/NEGLIGIBLE", "55%+",
          "of landlords say rental income is a minor or negligible share of total income — a side activity, not a profession.",
          value_color=ACCENT, value_size=30)
textbox(s, MARGIN, Inches(4.3), Inches(11), Inches(0.9), [
    {"text": "That's the design brief: lightweight, not sophisticated. A power-user portfolio tool would miss the actual market — the 6+ property segment is flat or shrinking as a share, not growing.", "size": 16, "color": INK}
])
source_line(s, "SOURCE: IW Köln, Deutschland.Immobilien Vermieterreport 2026", top=Inches(5.4))
footer(s, 4, TOTAL)

# ---------------------------------------------------------------- Slide 5
s = add_slide(prs)
eyebrow(s, "05", "THE PITCH")
title(s, "Your landlords are already your users. Give them the module they're missing.", size=27)
textbox(s, MARGIN, Inches(1.85), Inches(11), Inches(0.7), [
    {"text": "Chleo's platform already has private landlords as listing-side users. This isn't a new market — it's a second product for a customer she already has.", "size": 15, "color": MUTED}
])
steps = [
    ("01", "Upload", "Landlord uploads a rent receipt, EMI statement, Nebenkosten invoice, or lease."),
    ("02", "AI drafts", "The model extracts fields and rates its own confidence — never guesses silently."),
    ("03", "Landlord confirms", "Nothing saves until a human reviews it. Low confidence gets flagged, not hidden."),
    ("04", "Dashboard updates", "Portfolio profit & loss, occupancy, and market comps — always current."),
]
x = MARGIN
w = Inches(2.75)
for n, h3, p in steps:
    box = slide = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.8), w, Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = PAPER_RAISED
    box.line.color.rgb = RULE_STRONG; box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.14); tf.margin_right = Inches(0.14)
    p0 = tf.paragraphs[0]; r0 = p0.add_run(); _set_run(r0, n, size=11, color=ACCENT, font=MONO_FONT); p0.space_after = Pt(4)
    p1 = tf.add_paragraph(); r1 = p1.add_run(); _set_run(r1, h3, size=16, color=INK, bold=True); p1.space_after = Pt(6)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); _set_run(r2, p, size=11.5, color=MUTED)
    x = Emu(x + w + Inches(0.15))
footer(s, 5, TOTAL)

# ---------------------------------------------------------------- Slide 6
s = add_slide(prs)
eyebrow(s, "06", "ANSWERING THE FEAR DIRECTLY")
title(s, "“AI isn't transparent” — here's what we built so that objection doesn't hold.", size=25)
bullet_list(s, MARGIN, Inches(2.0), Inches(7.0), Inches(3.5), [
    "We cite the failure rate, not just the success rate. AI Index 2026's MortgageTax benchmark: even frontier models don't reliably exceed 69.4% accuracy on real financial documents. We designed around that number, not a marketing claim.",
    "Every extraction is a draft, never a save. The confidence threshold is set at that same 0.70 ceiling — not a round number picked for looks.",
    "Every model call is traced in LangSmith — what it saw, what it returned, why it was or wasn't flagged. Inspectable, not a black box.",
], size=13.5, gap=14)
field_box(s, Inches(8.4), Inches(2.0), Inches(4.2), Inches(1.5), "SAMPLE: CLEAN DOCUMENT", "confidence 0.98", "", value_size=20)
badge(s, Inches(8.55), Inches(3.15), "confirmed", "confirmed")
field_box(s, Inches(8.4), Inches(3.75), Inches(4.2), Inches(1.5), "SAMPLE: GARBLED SCAN", "confidence 0.50", "", value_size=20)
badge(s, Inches(8.55), Inches(4.9), "needs review", "flagged")
source_line(s, "SOURCE: AI Index Report 2026 §2.5 (MortgageTax) · langsmith/run_monitoring_sample.py, actually executed", top=Inches(6.5))
footer(s, 6, TOTAL)

# ---------------------------------------------------------------- Slide 7
s = add_slide(prs)
eyebrow(s, "07", "WHAT THE DASHBOARD REVEALS")
title(s, "The number a landlord tracks isn't the number that matters.", size=29)
textbox(s, MARGIN, Inches(1.85), Inches(11.5), Inches(0.9), [
    {"text": "Rent minus mortgage looks fine on most of our sample portfolio. The true monthly cash flow — after management fees, maintenance reserve, insurance, property tax, and Nebenkosten shortfalls — tells a different story.", "size": 14.5, "color": INK}
])
rows = [
    ("Property", "City", "Rent − EMI (naive)", "True net cash flow", "Status"),
    ("P02", "Leipzig", "+€352/mo", "+€55/mo", "still positive"),
    ("P01", "Berlin", "+€51/mo", "−€230/mo", "hidden loss"),
    ("P03", "München", "−€56/mo", "−€341/mo", "hidden loss"),
]
tbl_shape = s.shapes.add_table(len(rows), 5, MARGIN, Inches(3.0), Inches(11.5), Inches(1.8))
table = tbl_shape.table
widths = [Inches(2.2), Inches(2.0), Inches(2.6), Inches(2.6), Inches(2.1)]
for i, w in enumerate(widths):
    table.columns[i].width = w
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER_RAISED if r == 0 else PAPER
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]
        run.font.name = MONO_FONT if r > 0 else DISPLAY_FONT
        run.font.size = Pt(12 if r == 0 else 13)
        run.font.bold = (r == 0)
        run.font.color.rgb = MUTED if r == 0 else (FLAGGED if val == "hidden loss" else (CONFIRMED if val == "still positive" else INK))
textbox(s, MARGIN, Inches(5.2), Inches(11), Inches(0.7), [
    {"text": "6 of the 8 sample properties run a real monthly loss once full costs are counted — even though most nominally “cover” the mortgage. That gap is the product's reason to exist.", "size": 14, "color": INK, "italic": True}
])
source_line(s, "SOURCE: data/synthetic_landlord_portfolio.csv + synthetic_monthly_transactions.csv (fabricated, per the brief's data-source constraint)", top=Inches(6.6))
footer(s, 7, TOTAL)

# ---------------------------------------------------------------- Slide 8
s = add_slide(prs)
eyebrow(s, "08", "COST & TIMELINE")
title(s, "The AI itself is the cheap part. Say so.", size=30)
field_box(s, MARGIN, Inches(2.0), Inches(5.3), Inches(1.5), "COST PER SHORT DOCUMENT", "~€0.005–0.01",
          "Rent receipt / Nebenkosten invoice — GPT-4.1, $2/$8 per 1M in/out tokens.", value_color=ACCENT, value_size=24)
field_box(s, MARGIN, Inches(3.65), Inches(5.3), Inches(1.5), "COST PER LEASE DOCUMENT", "~€0.03–0.08",
          "Multi-page rental contract.", value_color=ACCENT, value_size=24)
bullet_list(s, Inches(6.7), Inches(2.0), Inches(5.9), Inches(2.6), [
    "Pilot hardening — ~2 weeks: auth, error handling, a real draft-record store.",
    "Pilot — ~8–10 weeks with ~200 opted-in landlords from the existing user base.",
    "Decision point — ~3 months out, on real pilot data, not assumptions.",
], size=13.5, gap=12)
field_box(s, Inches(6.7), Inches(4.7), Inches(5.9), Inches(1.4), "FIXED MONTHLY TOOLING COST", "~€70–90",
          "Dominates the total until landlord count reaches the thousands.", value_size=22)
source_line(s, "SOURCE: cost_estimation/cost_analysis.md — cited to OpenAI, Tableau, LangSmith public pricing (Aug 2026)", top=Inches(6.5))
footer(s, 8, TOTAL)

# ---------------------------------------------------------------- Slide 9
s = add_slide(prs)
eyebrow(s, "09", "RISKS, NAMED — NOT BURIED")
title(s, "What could go wrong, and where we've deliberately drawn the line.", size=27)
risks = [
    (MARGIN, Inches(2.0), FLAGGED, "TECHNICAL", "Document extraction tops out around 70% accuracy on real financial documents — designed around with mandatory human review, not hidden."),
    (Inches(6.9), Inches(2.0), FLAGGED, "REGULATORY", "Uploaded documents are sensitive financial data — a real GDPR surface (legal basis, DPIA, third-party processor) for Round 2, flagged now."),
    (MARGIN, Inches(4.0), MUTED, "EXPLICITLY OUT OF SCOPE", "Tenant screening/approval decisions. Colorado's AI Act names housing decisions as discrimination-risk, alongside hiring and medical care."),
    (Inches(6.9), Inches(4.0), MUTED, "ADOPTION", "Our own customer (slide 04) is price-sensitive and low-tech-sophistication — the MVP stays lightweight on purpose."),
]
for x, y, barcolor, k, p in risks:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), Inches(1.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = barcolor; bar.line.fill.background()
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x + Inches(0.06)), y, Inches(5.4), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = PAPER_RAISED; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.14); tf.margin_right = Inches(0.14)
    p0 = tf.paragraphs[0]; r0 = p0.add_run(); _set_run(r0, k, size=10.5, color=barcolor, bold=True, font=MONO_FONT); p0.space_after = Pt(6)
    p1 = tf.add_paragraph(); r1 = p1.add_run(); _set_run(r1, p, size=12.5, color=INK)
footer(s, 9, TOTAL)

# ---------------------------------------------------------------- Slide 10
s = add_slide(prs)
eyebrow(s, "10", "IF THIS SURVIVES THE ROOM")
title(s, "Round 2 deepens exactly this — it doesn't pivot to something new.", size=27)
bullet_list(s, MARGIN, Inches(2.0), Inches(7.0), Inches(3.8), [
    "A working MVP: document upload → extraction → confirm → dashboard, running end to end, not just described.",
    "A real ROI case, built from pilot data instead of the illustrative numbers on slide 08.",
    "EU AI Act classification and GDPR documentation — including the honest “minimal risk, and here's why” argument this use case supports.",
    "A strategic plan: pilot → full rollout, with the market-valuation use case as the next phase, not an afterthought.",
], size=13.5, gap=13)
field_box(s, Inches(8.4), Inches(2.0), Inches(4.2), Inches(2.2), "WHAT DOES NOT CHANGE",
          "Real estate · private landlords · this exact use case",
          "Unless today's feedback says otherwise.", value_size=16)
footer(s, 10, TOTAL)

# ---------------------------------------------------------------- Slide 11
s = add_slide(prs)
eyebrow(s, "11", "THE ASK")
title(s, "Tell us if this is the room we should be building in.", size=29)
field_box(s, MARGIN, Inches(2.1), Inches(11.5), Inches(1.2), "WE'RE ASKING FOR",
          "Feedback on the sector, the use case, and whether the transparency story actually lands — not just applause.", value_size=17)
bullet_list(s, MARGIN, Inches(3.7), Inches(11.5), Inches(2.4), [
    "Did the Hidden Cost Gap (slide 07) feel like a real insight, or an obvious one?",
    "Does “minimal, lightweight, for non-professional landlords” read as a strength or as too small a bet?",
    "Keep this industry and use case, or should we change direction before Round 2?",
], size=14.5, gap=12)
footer(s, 11, TOTAL)

# ---------------------------------------------------------------- Slide 12
s = add_slide(prs)
eyebrow(s, "12", "RECORDED AFTER THIS PRESENTATION")
title(s, "Keep or change — written down honestly, not decided in advance.", size=27)
field_box(s, MARGIN, Inches(2.1), Inches(5.6), Inches(1.9), "DECISION",
          "Pending — see feedback/round1_decision.md",
          "Intentionally not pre-filled. The feedback from this room determines it.", value_size=16)
textbox(s, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.0), [
    {"text": ("Changing industry or use case after this presentation costs nothing per "
              "the brief — and pretending a decision was already made before the room "
              "spoke would undercut the entire transparency pitch this deck just made."),
     "size": 14, "color": INK}
])
footer(s, 12, TOTAL)

# ---------------------------------------------------------------- Slide 13
s = add_slide(prs)
eyebrow(s, "13", "APPENDIX — MONETIZATION")
title(s, "Would €2 a month even be worth charging?", size=32)
textbox(s, MARGIN, Inches(1.85), Inches(11.5), Inches(1.3), [
    {"text": ("Assume 150,000 new rental properties enter the German market each year. "
              "That number isn't picked out of thin air — it independently matches our "
              "own cited landlord-growth data: private-landlord households grew by +1.8M "
              "over 2010–2022 (slide 03), which averages to almost exactly 150,000/year."),
     "size": 15, "color": INK}
])
add_table(
    s, MARGIN, Inches(3.35), Inches(11.5), Inches(1.9),
    ["Market capture", "Landlords", "Revenue / month", "Revenue / year"],
    [
        ["5%", "7,500", "€15,000", "€180,000"],
        ["10%", "15,000", "€30,000", "€360,000"],
        ["25%", "37,500", "€75,000", "€900,000"],
        ["100% (ceiling)", "150,000", "€300,000", "€3,600,000"],
    ],
    [Inches(3.0), Inches(2.8), Inches(2.8), Inches(2.9)],
)
textbox(s, MARGIN, Inches(5.55), Inches(11.5), Inches(0.9), [
    {"text": ("Even a coffee-money price point aggregates into real revenue at this market "
              "size — the question isn't whether €2 is enough, it's whether the "
              "cost side stays trivial enough to make that irrelevant. Slide 14 checks that directly."),
     "size": 13.5, "color": MUTED, "italic": True}
])
source_line(s, "SOURCE: IW Köln SOEP v39+Zensus (see slide 03) · 150,000/year is illustrative market sizing, not a verified external forecast", top=Inches(6.7))
footer(s, 13, TOTAL)

# ---------------------------------------------------------------- Slide 14
s = add_slide(prs)
eyebrow(s, "14", "APPENDIX — UNIT ECONOMICS")
title(s, "The math holds even at the lowest defensible price.", size=28)
field_box(s, MARGIN, Inches(2.0), Inches(5.3), Inches(1.5), "VARIABLE COST PER LANDLORD", "~€0.02–0.04/mo",
          "AI extraction cost, blended — from cost_estimation/cost_analysis.md.", value_color=ACCENT, value_size=22)
field_box(s, MARGIN, Inches(3.65), Inches(5.3), Inches(1.7), "FIXED MONTHLY COST", "~€2,665–2,685",
          "Tableau+n8n tooling (~€70–90) plus 0.1 FTE maintenance at the same €150/hr rate already used in the cost estimate.", value_color=ACCENT, value_size=22)
field_box(s, Inches(6.7), Inches(2.0), Inches(5.9), Inches(1.5), "BREAKEVEN", "~1,340 landlords",
          "Under 1% of the 150,000/year addressable pool — inside the pilot scale already discussed.", value_size=24)
add_table(
    s, Inches(6.7), Inches(3.7), Inches(5.9), Inches(1.4),
    ["Landlords", "Revenue", "Cost", "Margin"],
    [
        ["2,000 (pilot)", "€4,000/mo", "€2,765/mo", "+31%"],
        ["15,000 (10%)", "€30,000/mo", "€3,285/mo", "+89%"],
    ],
    [Inches(1.7), Inches(1.4), Inches(1.4), Inches(1.4)],
)
textbox(s, MARGIN, Inches(5.55), Inches(11.5), Inches(1.0), [
    {"text": ("Even the pilot scale already discussed elsewhere in this deck is profitable "
              "at €2/month. Honest caveat: this excludes customer acquisition, support, "
              "and churn — it's a floor showing the cost side isn't the constraint, not a "
              "recommended final price."),
     "size": 13, "color": INK}
])
source_line(s, "SOURCE: cost_estimation/cost_analysis.md — unit costs cited to public pricing; breakeven/margin computed here, not previously published", top=Inches(6.7))
footer(s, 14, TOTAL)

# ---------------------------------------------------------------- Slide 15
s = add_slide(prs)
eyebrow(s, "15", "PROJECT PLAN")
title(s, "What ships when — priority follows the customer, not the demo.", size=28)
textbox(s, MARGIN, Inches(1.85), Inches(11.5), Inches(0.9), [
    {"text": ("Scoped as a MoSCoW backlog. The order isn't about what's technically "
              "impressive — it's about what the landlord from slide 04 (small portfolio, "
              "side income, low tech-sophistication) actually needs first."),
     "size": 14.5, "color": INK}
])
add_table(
    s, MARGIN, Inches(3.0), Inches(11.5), Inches(3.0),
    ["Priority", "Feature", "Status"],
    [
        ["Must have", "Manage the costs & profits",
         "Built - Round 1 core (extraction + dashboard); Round 2 hardens to production"],
        ["Should have", "Tax savings advisor bot", "New - Round 2/3 candidate"],
        ["Could have", "Automated rental contract generation",
         "Adjacent to lease-extraction already built - inverse direction"],
        ["Could have", "Tenant complaints management",
         "Adjacent to maintenance-request tracking already prototyped"],
        ["Nice to have", "Direct bank connections (rent verification, tenant reminders, "
         "auto statement pull)", "Already demoed end to end (Enable Banking Mock ASPSP) - "
         "lowest priority, most de-risked"],
    ],
    [Inches(1.6), Inches(3.4), Inches(6.5)],
)
textbox(s, MARGIN, Inches(6.15), Inches(11.5), Inches(1.0), [
    {"text": ("That last row is deliberate: technical readiness and customer priority "
              "aren't the same axis. Bank connections are the furthest along technically "
              "and the lowest priority for this customer - proof the roadmap follows the "
              "persona (slide 04), not the excitement of what got built."),
     "size": 13, "color": INK, "italic": True}
])
footer(s, 15, TOTAL)

out_path = "slides.pptx"
prs.save(out_path)
print(f"Saved {out_path} with {len(prs.slides.slides)} slides" if hasattr(prs.slides, 'slides') else f"Saved {out_path} with {len(prs.slides._sldIdLst)} slides")
